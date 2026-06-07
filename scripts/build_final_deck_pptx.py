# -*- coding: utf-8 -*-
"""Editable PowerPoint (.pptx) mirror of the final jury deck.
Same content, layout and palette as scripts/build_final_deck.py (PDF), but native
text boxes, tables and shapes so the deck can be edited in PowerPoint.
Run: uv run python scripts/build_final_deck_pptx.py  ->  trexCloud_Final_Sunum.pptx
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---------- data ----------
B = json.loads(Path("web/public/data/bundle.json").read_text(encoding="utf-8"))
DEP = B["pm_value"]["deployed"]
SYNC = B["crossmachine"]["synchronization"]
MACHINES = sorted(
    [m for m in B["machines"] if m["has_telemetry"] and m["down_h"] > 0],
    key=lambda m: -m["oee"],
)
RISK = B["fanuc"]["risk"].get("Makine 1", [])
THR = B["fanuc"]["meta"]["threshold"]

# ---------- palette ----------
GREEN = RGBColor(0x0F, 0x7A, 0x52)
GREEN_SOFT = RGBColor(0xEC, 0xF4, 0xEF)
INK = RGBColor(0x16, 0x20, 0x1C)
MUTED = RGBColor(0x6B, 0x77, 0x6F)
HAIR = RGBColor(0xD9, 0xDF, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC4, 0x3D, 0x3D)
AMBER = RGBColor(0xE0, 0x74, 0x1C)
GRAY = RGBColor(0x9A, 0xA8, 0xA1)
TRACK = RGBColor(0xEE, 0xF2, 0xF0)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
ML, CW = 0.85, 13.333 - 1.7


def tu(s):
    return s.replace("i", "İ").upper()


def tr(x):
    return f"{x}".replace(".", ",")


def slide():
    return prs.slides.add_slide(BLANK)


def box(sl, l, t, w, h, fill=None, line=None, lw=1.0):
    sp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    return sp


def text(sl, s, l, t, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, lead=1.12):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, tf.margin_right):
        pass
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = lead
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    return tb


def bullets(sl, items, l, t, w, size=13.5, gap=6):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(7.0 - t - 0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        p.space_after = Pt(gap)
        b = p.add_run()
        b.text = "•  "
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.name = FONT
        b.font.color.rgb = GREEN
        r = p.add_run()
        r.text = it
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = INK
    return tb


def base(section, title, n):
    sl = slide()
    box(sl, 0, 0, 0.14, 7.5, fill=GREEN)
    if section:
        text(sl, tu(section), ML, 0.48, CW, 0.3, 11.5, GREEN, bold=True)
    text(sl, title, ML, 0.82, CW, 0.7, 26, INK, bold=True)
    box(sl, ML, 1.92, CW, 0.012, fill=HAIR)
    text(sl, "trexCloud · Kestirimci OEE ve Kök Neden Analizi", ML, 7.05, 8, 0.3, 9, MUTED)
    text(sl, f"{n:02d}", 13.333 - 1.4, 7.05, 0.55, 0.3, 9, MUTED, align=PP_ALIGN.RIGHT)
    return sl


def numbox(sl, items, top, h=1.45, gap=0.4):
    n = len(items)
    bw = (CW - (n - 1) * gap) / n
    for i, (head, sub) in enumerate(items):
        x = ML + i * (bw + gap)
        last = i == n - 1
        box(sl, x, top, bw, h, fill=(GREEN_SOFT if last else WHITE),
            line=(GREEN if last else HAIR), lw=(1.5 if last else 1))
        text(sl, str(i + 1), x + 0.16, top + 0.12, 0.5, 0.3, 12, GREEN if last else MUTED, bold=True)
        text(sl, head, x + 0.16, top + 0.45, bw - 0.3, 0.5, 13, GREEN if last else INK, bold=True)
        text(sl, sub, x + 0.16, top + h - 0.5, bw - 0.3, 0.45, 9.5, MUTED)


def statcards(sl, stats, top, h=1.35):
    n = len(stats)
    gap = 0.25
    swd = (CW - (n - 1) * gap) / n
    for i, (v, l) in enumerate(stats):
        x = ML + i * (swd + gap)
        box(sl, x, top, swd, h, fill=WHITE, line=HAIR, lw=1.0)
        text(sl, v, x, top + 0.18, swd, 0.5, 26, GREEN, bold=True, align=PP_ALIGN.CENTER)
        text(sl, l, x + 0.1, top + 0.88, swd - 0.2, 0.4, 10.5, MUTED, align=PP_ALIGN.CENTER)


def table(sl, headers, rows, widths, top, row_h=0.44, highlight0=True, centers=()):
    """Manual table from rectangles for exact look and easy edit."""
    x0 = ML
    # header
    x = x0
    for j, htext in enumerate(headers):
        box(sl, x, top, widths[j], 0.42, fill=INK)
        text(sl, htext, x + 0.1, top + 0.09, widths[j] - 0.18, 0.3, 10.5, WHITE, bold=True)
        x += widths[j]
    for i, row in enumerate(rows):
        y = top + 0.42 + i * row_h
        hi = highlight0 and i == 0
        x = x0
        for j, val in enumerate(row):
            box(sl, x, y, widths[j], row_h, fill=(GREEN_SOFT if hi else WHITE), line=HAIR, lw=0.6)
            al = PP_ALIGN.CENTER if j in centers else PP_ALIGN.LEFT
            text(sl, val, x + 0.1, y + 0.1, widths[j] - 0.18, row_h - 0.1, 10,
                 GREEN if hi else INK, bold=hi, align=al)
            x += widths[j]


# ═══════════════════════════ 01 cover ═══════════════════════════
sl = slide()
box(sl, 0, 0, 0.14, 7.5, fill=GREEN)
text(sl, "TREXCLOUD · ULUDAĞ HACKATHON", ML, 2.3, CW, 0.3, 13, GREEN, bold=True)
text(sl, "Kestirimci OEE ve Kök Neden Analizi", ML, 2.75, CW, 0.9, 38, INK, bold=True)
text(sl, "Bir CNC ve lazer tesisi için uçtan uca duruş tahmini, kök neden ve What-If sistemi",
     ML, 4.4, CW, 0.4, 16, MUTED)
box(sl, ML, 5.2, 2.2, 0.025, fill=GREEN)
text(sl, "Veri dönemi Ağustos 2025 - Mayıs 2026  ·  12 makine  ·  yaklaşık 7,4 milyon telemetri kaydı",
     ML, 5.4, CW, 0.3, 12.5, MUTED)
text(sl, "trexCloud · Kestirimci OEE ve Kök Neden Analizi", ML, 7.05, 8, 0.3, 9, MUTED)

# ═══════════════════════════ 02 veri ═══════════════════════════
sl = base("Veri", "Önce veriyi dürüstçe envanterledik", 2)
bullets(sl, [
    "Tesiste 12 makine var (Fanuc, Mitsubishi, Nukon). Yaklaşık 7,4 milyon telemetri kaydı ve "
    "153 bin MES olayı işlendi. Tüm süreler milisaniye cinsindendir.",
    "OEE üç bileşenin çarpımıdır: Kullanılabilirlik (A) çarpı Performans (P) çarpı Kalite (Q).",
    "Veride hurda kaydı sıfırdır, bu yüzden Kalite her zaman 1 çıkar. Üretim sayımı sıfır olan "
    "günlerde Performans 0 olur.",
    "Dokümanın işaret ettiği kanıt sinyalleri (servo sıcaklığı, yük) pratikte boştur. Modeli "
    "gerçekten akan sinyallerin üzerine kurduk: cycle_time, run_state, run_time, axis_position.",
    "En büyük duruş kaynağı bir makine arızası değil, System Offline yani bağlantı kesintisidir. "
    "Bunu ayrı tuttuk; çözümü bakım değil IT ekibindedir ve makine OEE'sine yazılmaz.",
], ML, 2.25, CW, 13.5, gap=10)

# ═══════════════════════════ 03 yöntem ═══════════════════════════
sl = base("Yöntem", "Sinyalden öneriye uzanan tek hat", 3)
text(sl, "Sistem beş aşamadan oluşur. Her aşama bir öncekinin çıktısını kullanır.", ML, 2.2, CW, 0.3, 14.5, INK)
numbox(sl, [
    ("Kanonik sinyal", "Farklı markaların sinyalleri ortak rollere eşlenir"),
    ("Baseline sapma", "Her sinyalin normalinden sapması ölçülür"),
    ("Kök neden", "Alarmlar nedensel zincire dizilir"),
    ("Duruş tahmini", "HistGBDT modeli riski önceden işaretler"),
    ("What-If ve OEE", "Düzeltmenin OEE kazanımı sayısallaşır"),
], 2.85)
bullets(sl, [
    "Marka bağımsız rol haritası makineler arası karşılaştırmayı mümkün kılar. Örneğin Fanuc "
    "IsNotRunning sinyali ile Mitsubishi RUN_STATUS_START sinyali aynı role eşlenir.",
    "Her makinenin sinyalleri kendi içinde sağlam istatistiklerle (medyan, çeyrekler açıklığı) "
    "normalize edilir, böylece farklı marka ölçekleri ortak bir uzaya taşınır.",
    "Model eğitimi sızıntıya karşı korumalıdır. Zaman sırasına göre bölünür, yalnızca geçmiş "
    "öznitelikler kullanılır ve normalizasyon istatistikleri yalnızca eğitim verisinden hesaplanır.",
], ML, 4.7, CW, 12.5, gap=6)

# ═══════════════════════════ 04 filo (chart) ═══════════════════════════
sl = base("Filo Durumu", "Makine bazında OEE ve ana kayıp kaynağı", 4)
text(sl, "Telemetrisi olan ve duruşu kaydedilen yedi makine, OEE değerine göre sıralanmıştır. "
     "Yeşil Fanuc, turuncu Mitsubishi makineleri gösterir.", ML, 2.2, CW, 0.5, 13.5, INK)
# bars
bx, bt, bw_total, bh = ML, 2.95, 6.0, 3.3
rh = bh / len(MACHINES)
for i, m in enumerate(MACHINES):
    y = bt + i * rh
    col = GREEN if m["regime"].startswith("Fanuc") else AMBER
    text(sl, m["name"].replace("Makine", "M"), bx, y + rh * 0.22, 0.9, 0.3, 10.5, INK)
    tx, tw = bx + 0.95, bw_total - 2.3
    box(sl, tx, y + rh * 0.18, tw, rh * 0.5, fill=TRACK)
    box(sl, tx, y + rh * 0.18, max(0.03, tw * m["oee"]), rh * 0.5, fill=col)
    text(sl, f"{m['oee']*100:.1f}%", tx + tw + 0.1, y + rh * 0.2, 1.1, 0.3, 10.5, INK, bold=True)
box(sl, 7.4, 2.95, 5.05, 3.3, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(sl, "Ne görüyoruz?", 7.62, 3.15, 4.7, 0.3, 12, GREEN, bold=True)
bullets(sl, [
    "OEE değerleri yüzde 0 ile yüzde 74 arasında dağılır. En iyi makine bile yüzde 75'in altındadır, "
    "yani iyileştirme alanı geniştir.",
    "Makine 8'in OEE değeri 0'dır çünkü o dönemde üretim kaydı yoktur. Duruşu gerçektir, fakat "
    "Performans çarpanı 0 olduğu için OEE sıfırlanır.",
    "Tüm makinelerde OEE pratikte Kullanılabilirlik tarafından belirlenir.",
], 7.62, 3.5, 4.55, 10.5, gap=6)

# ═══════════════════════════ 05 kök neden ═══════════════════════════
sl = base("Kök Neden", "Baseline sapma ve alarm zinciri", 5)
text(sl, "Örnek olay: Makine 1, 12 Ocak 2026 saat 04:47.", ML, 2.2, CW, 0.3, 14, MUTED)
box(sl, ML, 2.85, 3.6, 0.9, fill=GREEN_SOFT, line=GREEN, lw=1.5)
text(sl, "AIR PRESSURE FAILED", ML, 3.0, 3.6, 0.3, 13.5, GREEN, bold=True, align=PP_ALIGN.CENTER)
text(sl, "kök neden", ML, 3.35, 3.6, 0.3, 11, GREEN, align=PP_ALIGN.CENTER)
box(sl, ML + 4.35, 2.85, 3.6, 0.9, fill=WHITE, line=HAIR, lw=1.0)
text(sl, "Z-AXIS ZERO RETURN", ML + 4.35, 3.0, 3.6, 0.3, 13.5, INK, bold=True, align=PP_ALIGN.CENTER)
text(sl, "bu alarmın sonucu", ML + 4.35, 3.35, 3.6, 0.3, 11, MUTED, align=PP_ALIGN.CENTER)
bullets(sl, [
    "Olay anında birden fazla sinyal baseline değerinden sapar. run_state sinyali eksi 7,7 standart "
    "sapmaya iner, yani makine durur. cycle_time eksi 0,8 standart sapma gösterir.",
    "Aynı anda gelen alarmlar dizine göre değil, fiziksel nedensel önceliğe göre sıralanır. "
    "Kök neden hava basıncıdır, Z ekseni hatası onun sonucudur.",
    "Sistemin çıktısı yalnızca bir alarm değildir. Kanıtı, sıralı hipotezleri ve uygulanabilir "
    "bakım aksiyonunu içeren bir kök neden kartıdır.",
], ML, 4.1, CW, 13.5, gap=8)

# ═══════════════════════════ 06 model rejim (table) ═══════════════════════════
sl = base("Tahmin", "Duruş tahmin modeli ve eğitilen alternatifler", 6)
statcards(sl, [("0,73", "ROC-AUC · Fanuc"), ("2,38×", "Lift · Fanuc"),
               (f"%{DEP['recall']*100:.0f}", "Recall · held-out"), ("1,2×", "Lift · Mitsubishi")], 2.2, h=1.2)
table(sl,
      ["Eğitim verisi", "Test", "ROC", "Lift", "Recall", "Not"],
      [["Fanuc (ayarlanmış)", "Fanuc", "0,73", "2,38×", "%69", "Dağıtılan model"],
       ["Fanuc (temel)", "Fanuc", "0,68", "2,0×", "%59", "Izgara aramadan önce"],
       ["Birleşik + z-norm", "Fanuc", "0,72", "2,33×", "%48", "Marka transferi katkı vermedi"],
       ["Mitsubishi", "Mitsubishi", "0,61", "1,2×", "%97", "Taban %53, yüksek recall yanıltıcı"]],
      [3.0, 1.7, 1.0, 1.0, 1.1, 3.83], 3.65, centers=(2, 3, 4))
text(sl, "Fanuc hücresi {1,2,3,5,9} tek gerçek tahmin edilebilir gruptur. Mitsubishi makineleri zamanın "
     "yüzde 53'ünde zaten durur, bu yüzden yüksek recall değeri bilgi taşımaz ve lift yalnızca 1,2 katıdır.",
     ML, 6.25, CW, 0.5, 12, MUTED)

# ═══════════════════════════ 07 algoritma (table) ═══════════════════════════
sl = base("Tahmin", "Hangi algoritmayı neden seçtik", 7)
text(sl, "Aynı held-out gelecekte, aynı öznitelik kümesiyle dört model ailesini ve iki saçma tabanı "
     "yarıştırdık. Karşılaştırma birleşik veride yapıldı (taban oran yüzde 25).", ML, 2.2, CW, 0.5, 13.5, INK)
table(sl,
      ["Model ailesi", "ROC", "Lift", "Not"],
      [["HistGBDT (gradyan artırma)", "0,76", "2,09", "Seçilen aile, eksik değeri yerel işler"],
       ["Random Forest", "0,75", "2,00", "Yakın ikinci"],
       ["Lojistik Regresyon", "0,74", "1,98", "Doğrusal taban"],
       ["MLP (sinir ağı, 64-24)", "0,68", "1,72", "Tablo veride en zayıf öğrenen"],
       ["Denetimsiz AD skoru", "0,51", "1,07", "Tahminci olarak tesadüf seviyesi"],
       ["Her zaman pozitif (taban)", "0,50", "1,00", "Anlamsız referans"]],
      [3.7, 1.2, 1.2, 5.53], 2.95, centers=(1, 2))
box(sl, ML, 6.05, CW, 0.55, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(sl, "Neden HistGBDT: en yüksek PR-AUC ve ROC. Eksik değeri kendisi işler, ölçekleme istemez ve "
     "doğrusal olmayan eşikleri yakalar. CNN ile transformer yalnızca denetimsiz anomali "
     "otokodlayıcısında denendi, tahmin için avantaj sağlamadı.",
     ML + 0.22, 6.15, CW - 0.44, 0.4, 10, INK)

# ═══════════════════════════ 08 risk timeline (chart) ═══════════════════════════
sl = base("Tahmin", "Modelin ürettiği risk zaman çizgisi", 8)
text(sl, "Aşağıdaki eğri, Makine 1 için modelin held-out dönemde ürettiği duruş riskidir. Bu, hiçbir "
     "müdahale yapılmadığında beklenen riski gösterir.", ML, 2.2, CW, 0.5, 13.5, INK)
cx, ct, cw, ch = ML, 2.9, CW, 2.35
box(sl, cx, ct, cw, ch, fill=WHITE, line=HAIR, lw=1.0)
pad = 0.18
px, py, pw, ph = cx + pad + 0.3, ct + pad, cw - 2 * pad - 0.4, ch - 2 * pad - 0.2
for gv in (0.0, 0.5, 1.0):
    gy = py + ph * (1 - gv)
    ln = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(gy), Inches(px + pw), Inches(gy))
    ln.line.color.rgb = HAIR
    ln.line.width = Pt(0.4)
    text(sl, tr(f"{gv:.1f}"), cx + pad - 0.1, gy - 0.1, 0.35, 0.2, 8, MUTED, align=PP_ALIGN.RIGHT)
ty = py + ph * (1 - THR)
thln = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(ty), Inches(px + pw), Inches(ty))
thln.line.color.rgb = RED
thln.line.width = Pt(0.9)
thln.line.dash_style = 2  # dash
text(sl, "alarm eşiği", px + pw - 1.0, ty - 0.18, 1.0, 0.2, 8, RED, align=PP_ALIGN.RIGHT)
pts = RISK[:: max(1, len(RISK) // 150)] if RISK else []
for j in range(1, len(pts)):
    x1 = px + pw * ((j - 1) / (len(pts) - 1))
    x2 = px + pw * (j / (len(pts) - 1))
    y1 = py + ph * (1 - max(0.0, min(1.0, pts[j - 1]["r"])))
    y2 = py + ph * (1 - max(0.0, min(1.0, pts[j]["r"])))
    seg = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    seg.line.color.rgb = GREEN
    seg.line.width = Pt(1.0)
bullets(sl, [
    "Yeşil eğri modelin tahmin ettiği duruş riskidir, 0 ile 1 arasında değişir.",
    "Kırmızı kesikli çizgi alarm eşiğidir. Risk bu çizgiyi aştığında bir uyarı üretilir.",
    "Dönemin başında risk düşük ve düzdür. İlerleyen haftalarda risk sık sık eşiğin üzerine çıkar; "
    "model bu noktalarda yaklaşan duruşları işaretler.",
], ML, 5.45, CW, 12, gap=6)

# ═══════════════════════════ 09 bridge ═══════════════════════════
sl = base("OEE Kazanımı", "Tahmin OEE kazanımına nasıl bağlanıyor", 9)
text(sl, "Tahminin OEE değeriyle bağı, ölçülmüş üç gerçeğe dayanır.", ML, 2.2, CW, 0.3, 14, INK)
numbox(sl, [
    (f"%{DEP['recall']*100:.0f} recall", "önemli duruşların yakalanan oranı"),
    (f"{DEP['caught_downtime_h']:.0f} saat", "yakalanan duruşların toplam süresi"),
    ("Önceden işaretlenir", "körlemesine değil, önceden bilinerek ele alınır"),
], 2.85, h=1.45)
bullets(sl, [
    "Model held-out dönemde önemli duruşların yüzde 69'unu yakaladı. Bu duruşların içindeki toplam "
    "duruş süresi 2062 saattir; bu, Fanuc hücresinin tüm plansız duruş süresinin yaklaşık yüzde 75'idir.",
    "Yakalanan duruşlar artık körlemesine yaşanmaz, önceden bilinir. Böylece duruş süresi ele "
    "alınabilir bir hedefe dönüşür.",
    "Bu sürenin azaltılmasının OEE etkisi What-If senaryolarında somutlaşır. En büyük plansız duruş "
    "kalemini yüzde 30 azaltmak Makine 1'de OEE'yi 18, Makine 5'te 15 puan yükseltir.",
], ML, 4.75, CW, 13.5, gap=8)

# ═══════════════════════════ 10 örüntü (chart) ═══════════════════════════
sl = base("Örüntü", "Makineler arası örüntü, senkron duruşlar", 10)
text(sl, "Soru şu: iki veya daha fazla makine aynı saatte plansız durduğunda, bu durum şansın ve ortak "
     "vardiya programının ötesinde bir bağ mıdır?", ML, 2.2, CW, 0.5, 14, INK)
sbars = [("Gözlemlenen", SYNC["observed_co_stop_hours"], GREEN, "gerçekte yaşanan eş zamanlı duruş"),
         ("Beklenen, vardiya ritmi", round(SYNC["daily"]["exp"]), AMBER, "yalnız ortak saat düzeni açıklasaydı"),
         ("Beklenen, rastgele", round(SYNC["free"]["exp"]), GRAY, "tamamen tesadüf olsaydı")]
for i, (lab, v, col, note) in enumerate(sbars):
    y = 3.05 + i * 0.62
    text(sl, lab, ML, y, 3.0, 0.3, 12, INK, bold=(i == 0))
    bw_ = (v / 760) * 5.4
    box(sl, ML + 3.05, y, max(0.05, bw_), 0.38, fill=col)
    text(sl, f"{v} saat", ML + 3.2 + bw_, y + 0.03, 1.2, 0.3, 11.5, INK, bold=True)
    note_x = ML + 3.2 + bw_ + 1.05
    text(sl, note, note_x, y + 0.04, max(1.5, 13.05 - note_x), 0.3, 10, MUTED)
text(sl, f"İstatistik sonucu: z eşittir {tr(SYNC['daily']['z'])}, p küçüktür 0,001. Senkronizasyon ortak "
     "vardiya programıyla açıklanamıyor.", ML, 5.0, CW, 0.4, 13.5, INK)
box(sl, ML, 5.5, CW, 0.95, fill=GREEN_SOFT, line=GREEN, lw=1.2)
text(sl, "İçgörü: 708 saatlik eş zamanlı duruşun 564 saati ortak vardiya ile açıklanır. Geri kalan "
     "yaklaşık 144 saat vardiyayla açıklanamayan gerçek bir bağdır. Bu, duruşların bir bölümünün "
     "makineye özel değil sistemik olduğunu gösterir. Ortak bir kök nedeni (örneğin paylaşılan besleme "
     "veya altyapı) düzeltmek, aynı anda birden çok makinenin OEE değerini yükseltebilir.",
     ML + 0.22, 5.62, CW - 0.44, 0.8, 11, INK)

# ═══════════════════════════ 11 senaryolar (table) ═══════════════════════════
sl = base("Senaryolar", "What-If senaryoları, fiziksel kazanım", 11)
text(sl, "Her senaryo ham OEE bileşenlerinden aynı formülle yeniden hesaplanır. Değerler fiziksel "
     "kazanımdır; para birimi içermez.", ML, 2.2, CW, 0.5, 13.5, INK)
table(sl,
      ["Senaryo", "Kaldıraç", "ΔOEE", "Kazanılan süre"],
      [["S1 · Makine 1 · en büyük plansız kalem -%30", "Kullanılabilirlik", "+18,0 pp", "625 saat çalışma"],
       ["S2 · Makine 5 · en büyük plansız kalem -%30", "Kullanılabilirlik", "+15,3 pp", "420 saat çalışma"],
       ["S3 · Makine 3 · çevrim süresi +%10", "Performans", "0,0 pp", "veri P kolunu desteklemiyor"],
       ["S4 · Tesis · bağlantı kesintisi giderilir", "Bağlantı / IT", "makine dışı", "663 saat bağlantı"]],
      [5.6, 1.6, 1.3, 2.13], 2.78, row_h=0.5, centers=(2,))
bullets(sl, [
    "S1 ve S2 Kullanılabilirlik kaldıracıdır ve gerçek kazanım üretir. En büyük plansız duruş kalemini "
    "azaltmak iki farklı makinede de OEE'yi belirgin biçimde yükseltir.",
    "S3 Performans kaldıracını test eder ve ölçülen sonuç sıfırdır. Üretim yapılan 777 makine gününün "
    "605'inde Performans tam 1, 172'sinde 0'dır ve arada hiçbir gün yoktur. Bu veride çevrim "
    "iyileştirmesi OEE'yi hareket ettiremez.",
    "S4 bağlantı kesintisinin giderilmesidir. 663 saatlik bağlantı görünürlüğü geri kazanılır, fakat "
    "bu bir IT aksiyonudur ve makine OEE'sine yazılmaz. Kalite kaldıracı da hurda kaydı sıfır olduğu "
    "için inerttir.",
], ML, 5.45, CW, 10.5, gap=5)

# ═══════════════════════════ 12 sonuç ═══════════════════════════
sl = base("Sonuç", "Tahmin eder, açıklar, sayısallaştırır", 12)
statcards(sl, [("2,38×", "Fanuc duruş tahmini lift"),
               ("2062 saat", "önceden yakalanan duruş süresi"),
               ("z = 5,16", "senkron duruş anlamlılığı"),
               ("+18 puan", "en iyi What-If OEE kazanımı")], 2.45)
bullets(sl, [
    "Sistem tek bir hatta üç işi birlikte yapar: Fanuc hücresinde duruşu önceden tahmin eder, kök "
    "nedeni alarm zinciriyle açıklar ve düzeltmenin OEE kazanımını sayısallaştırır.",
    "Her iddia bir null model ile sınanmıştır. Olumsuz bulguları gizlemedik, çünkü güven dürüstlükten "
    "gelir.",
    "Canlı arayüz tahminden kök nedene ve What-If analizine uzanan akışı uçtan uca gösterir.",
], ML, 4.5, CW, 13.5, gap=10)

prs.save("trexCloud_Final_Sunum.pptx")
print(f"kaydedildi: trexCloud_Final_Sunum.pptx  ({len(prs.slides)} sayfa)")
