# -*- coding: utf-8 -*-
"""Build the presentation deck (16:9 .pptx) — professional, restrained, consistent master.

Every slide shares the same fixed-position furniture: a thin green accent strip, a section
label, the title baseline, a hairline rule, and the footer. Content sits in one consistent
grid. Speaker script goes into the notes pane so slides stay clean. Turkish, comma decimals.

Run: uv run python scripts/build_slides.py
"""
import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

PM = json.loads(Path("analysis/artifacts/pm_value.json").read_text(encoding="utf-8"))
SCENARIOS = json.loads(Path("analysis/artifacts/scenarios.json").read_text(encoding="utf-8"))["rows"]
DEPLOYED = PM["deployed"]
ECON = PM["sensitivity"]["economic_optimum"]
DASH = Path("analysis/reports/dashboard")

OUT = os.environ.get("TREX_SLIDES_OUT", "trexCloud_Sunum.pptx")
FONT = "Helvetica Neue"

INK = RGBColor(0x1A, 0x1F, 0x1D)
GREEN = RGBColor(0x12, 0x71, 0x45)       # deep corporate green
GREEN_SOFT = RGBColor(0xEC, 0xF4, 0xEF)
MUTED = RGBColor(0x6B, 0x77, 0x71)
HAIR = RGBColor(0xD9, 0xDF, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRONZE = RGBColor(0xA9, 0x77, 0x3B)
SILVER = RGBColor(0x8A, 0x96, 0x9E)
GOLD = RGBColor(0xC9, 0x8B, 0x2A)
PLAT = RGBColor(0x2F, 0x8E, 0x86)

SW, SH = 13.333, 7.5
ML = 0.85                 # left margin
CW = SW - ML - 0.85       # content width
TITLE_Y = 0.92
RULE_Y = 1.92
BODY_Y = 2.28
FOOT_Y = 7.04

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def _set(run, size, color=INK, bold=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color


def rich(para, text, size, color=INK, bold=False):
    """split on '*' to toggle bold emphasis."""
    for i, seg in enumerate(text.split("*")):
        if seg == "":
            continue
        r = para.add_run(); r.text = seg
        _set(r, size, color, bold or (i % 2 == 1))


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def text(slide, x, y, w, h, lines, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, gap=8, lh=1.12):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(gap); p.line_spacing = lh
        rich(p, ln, size, color, bold)
    return tb


def bullets(slide, items, y=BODY_Y, x=ML, w=CW, size=16, gap=11):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(SH - y - 0.9))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.16
        m = p.add_run(); m.text = "—  "; _set(m, size, GREEN, True)
        rich(p, it, size, INK)
    return tb


def base(slide, section, title, n):
    box(slide, 0, 0, 0.14, SH, fill=GREEN)                       # accent strip
    if section:
        text(slide, ML, 0.52, CW, 0.3, section.upper(), 11.5, GREEN, True)
    text(slide, ML, TITLE_Y, CW, 0.95, title, 29, INK, True, lh=1.02)
    box(slide, ML, RULE_Y, CW, 0.012, fill=HAIR)                 # hairline rule
    text(slide, ML, FOOT_Y, CW - 0.6, 0.3, "trexCloud · Öngörücü OEE & Kök-Neden Analizi",
         9, MUTED)
    text(slide, SW - 1.4, FOOT_Y, 0.55, 0.3, f"{n:02d}", 9, MUTED, align=PP_ALIGN.RIGHT)


def notes(slide, t):
    slide.notes_slide.notes_text_frame.text = t


def picture_fit(slide, path, x, y, w, h):
    """Place an image inside a box without stretching it."""
    path = Path(path)
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(
        str(path),
        Inches(x + (w - pw) / 2),
        Inches(y + (h - ph) / 2),
        width=Inches(pw),
        height=Inches(ph),
    )


def new(section, title, n):
    s = prs.slides.add_slide(BLANK)
    base(s, section, title, n)
    return s


# ───────────────────────── 01 · cover ─────────────────────────
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 0.14, SH, fill=GREEN)
text(s, ML, 2.35, CW, 0.4, "TREXCLOUD HACKATHON", 13, GREEN, True)
text(s, ML, 2.95, CW, 1.6, "Öngörücü OEE & Kök-Neden Analizi", 40, INK, True, lh=1.0)
text(s, ML, 4.35, CW, 0.6, "Bir CNC + lazer tesisi için uçtan uca tahmin, kök-neden ve What-If sistemi",
     16, MUTED)
box(s, ML, 5.15, 2.2, 0.02, fill=GREEN)
text(s, ML, 5.35, CW, 0.4, "Veri dönemi: Ağustos 2025 – Mayıs 2026  ·  12 makine", 12.5, MUTED)
text(s, ML, FOOT_Y, CW, 0.3, "trexCloud · Öngörücü OEE & Kök-Neden Analizi", 9, MUTED)
notes(s, "Kısa ve iddialı giriş: bir CNC + lazer tesisinin verisinden uçtan uca bir öngörücü "
         "OEE ve kök-neden sistemi kurduk. Hedefimiz değerlendirmenin en üst iki seviyesi: Altın ve Platin.")

# ───────────────────────── 02 · yaklaşım / madalya ─────────────────────────
s = new("Yaklaşım", "Hedef: değerlendirmenin en üst iki seviyesi", 2)
text(s, ML, BODY_Y, CW, 0.4,
     "Değerlendirme dört seviyeli. Bronz ve Gümüş bizde temel; odağımız Altın ve Platin.", 16, INK)
cards = [("BRONZ", BRONZE, "Alarm ve duruşları zaman aralığında listele", False),
         ("GÜMÜŞ", SILVER, "Alarmları duruşlarla eşleştir, Pareto çıkar", False),
         ("ALTIN", GOLD, "Baseline sapma + çok-sinyal → nedensellik zinciri", True),
         ("PLATİN", PLAT, "Çapraz-makine örüntü + ΔOEE + finansal öneri", True)]
cw = (CW - 3 * 0.3) / 4
for i, (name, col, desc, hi) in enumerate(cards):
    x = ML + i * (cw + 0.3)
    box(s, x, 3.1, cw, 2.7, fill=(GREEN_SOFT if hi else WHITE), line=(GREEN if hi else HAIR),
        lw=(1.5 if hi else 1.0))
    box(s, x, 3.1, cw, 0.12, fill=col)
    text(s, x + 0.18, 3.45, cw - 0.36, 0.4, name, 14, INK, True)
    text(s, x + 0.18, 3.95, cw - 0.36, 1.7, desc, 12.5, (INK if hi else MUTED), lh=1.18)
    if hi:
        text(s, x + 0.18, 5.35, cw - 0.36, 0.3, "HEDEF", 9.5, GREEN, True)
notes(s, "Jüriye yol haritasını ver: her rozetin ne istediğini bir cümlede söyle. Alttaki iki "
         "seviyeyi eksiksiz yapıyoruz; sunumun ağırlığı Altın ve Platin'de olacak.")

# ───────────────────────── 03 · veri gerçeği ─────────────────────────
s = new("Veri", "Önce veriyi dürüstçe envanterledik", 3)
bullets(s, [
    "12 makine, yaklaşık *7,4 milyon* telemetri kaydı; cp1254 kodlama, UTC, milisaniye.",
    "Dokümanın vaat ettiği kanıt sinyalleri (*servo sıcaklık, path-load*) pratikte boş: "
    "en yoğun makinede bile *0 / 28 satır*.",
    "Gerçekte akan sinyaller: *cycle_time, run_state, run_time, axis_position*. "
    "Nightwatch–MES eşleşmesi *162 / 162 (%100)*.",
], y=BODY_Y)
box(s, ML, 5.2, CW, 1.05, fill=GREEN_SOFT, line=GREEN, lw=1.2)
text(s, ML + 0.25, 5.45, CW - 0.5, 0.6,
     "Modeli dokümanın değil, gerçekten akan sinyallerin üzerine kurduk. "
     "Jüri veriyi iyi bilir; dürüst envanter güven kazandırır.", 14, INK, lh=1.2)
notes(s, "En güçlü farkımız: dokümanın vaat ettiği zengin sinyallerin akmadığını tespit ettik "
         "ve modeli gerçekte akan sinyallere kurduk. Bu, jüri gözünde güvenilirlik demektir.")

# ───────────────────────── 04 · mimari ─────────────────────────
s = new("Mimari", "Uçtan uca hat", 4)
steps = ["Kanonik sinyal\nkatmanı", "Baseline sapma\n(AD)", "Kök-neden\n(RCA)",
         "What-If / OEE", "Tahmin"]
bw = (CW - 4 * 0.45) / 5
for i, st in enumerate(steps):
    x = ML + i * (bw + 0.45)
    last = i == len(steps) - 1
    box(s, x, 3.0, bw, 1.5, fill=(GREEN_SOFT if last else WHITE),
        line=(GREEN if last else HAIR), lw=(1.5 if last else 1.0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + 0.1, 3.0, bw - 0.2, 1.5, st.split("\n"), 12.5,
         (GREEN if last else INK), True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, gap=1, lh=1.05)
    if not last:
        ar = box(s, x + bw + 0.08, 3.55, 0.3, 0.4, fill=HAIR, shape=MSO_SHAPE.RIGHT_ARROW)
        ar.line.fill.background()
text(s, ML, 5.2, CW, 0.5,
     "Vendor-bağımsız rol haritası çapraz-makineyi mümkün kılar: Fanuc IsNotRunning ile "
     "Mitsubishi RUN_STATUS_START aynı role eşlenir.", 14, MUTED, lh=1.2)
text(s, ML, 5.95, CW, 0.5,
     "Rejim ayrımı:  *Fanuc {1,2,3,5,9}*   ·   *Mitsubishi {7,8}*   ·   telemetrisiz {4,6,10, TurboCut, ARES}",
     13.5, INK)
notes(s, "Sistemi bir resimde anlat. Kanonik rol katmanı sayesinde farklı markaların sinyallerini "
         "ortak bir dile çeviriyoruz; tahmini ve çapraz-makine analizini bu mümkün kılıyor.")

# ───────────────────────── 05 · bronz + gümüş ─────────────────────────
s = new("Bronz · Gümüş", "Temel: olay akışı, eşleştirme, Pareto", 5)
bullets(s, [
    "Alarm, duruş, bağlantı kesintisi ve sapma olayları tek bir akışta birleşik.",
    "Alarmlar duruşlara zaman damgasıyla eşleştirildi (ileri yönlü asof).",
    "Pareto, duruşu *arıza* ve *bağlantı* olarak ayırır — yatırım kararını bu ayrım belirler.",
], y=BODY_Y)
box(s, ML, 5.2, CW, 1.05, fill=GREEN_SOFT, line=GREEN, lw=1.2)
text(s, ML + 0.25, 5.45, CW - 0.5, 0.6,
     "En büyük duruş kaynağı bir makine arızası değil, bağlantı kesintisidir (System Offline). "
     "Bunu ayırmazsanız yanlış yere yatırım yaparsınız.", 14, INK, lh=1.2)
notes(s, "Hızlı geç. Kritik ayrım: en büyük duruş kalemi aslında bir IT/ağ sorunu. Bunu ayırmak "
         "doğru aksiyonun ön koşulu.")

# ───────────────────────── 06 · ALTIN ─────────────────────────
s = new("Altın", "Baseline sapma + çok-sinyal → nedensellik", 6)
text(s, ML, BODY_Y, CW, 0.35, "Örnek olay:  Makine 1  ·  12 Ocak 2026, 04:47", 14, MUTED)
# cascade
c1 = box(s, ML, 2.95, 3.6, 0.9, fill=GREEN_SOFT, line=GREEN, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, ML, 2.95, 3.6, 0.9, ["AIR PRESSURE FAILED", "kök neden"], 13.5, GREEN, True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, gap=2, lh=1.0)
ar = box(s, ML + 3.75, 3.22, 0.45, 0.45, fill=HAIR, shape=MSO_SHAPE.RIGHT_ARROW); ar.line.fill.background()
box(s, ML + 4.35, 2.95, 3.6, 0.9, fill=WHITE, line=HAIR, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, ML + 4.35, 2.95, 3.6, 0.9, ["Z-AXIS ZERO RETURN", "sonuç"], 13.5, INK, True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, gap=2, lh=1.0)
bullets(s, [
    "Baseline sapma imzası (robust-z):  *run_state −7,7σ* (makine duruyor),  cycle_time −0,8σ.",
    "Alarmlar dizine göre değil *nedensel önceliğe* göre sıralanır: kök neden hava basıncı, "
    "Z-ekseni hatası sonuçtur.",
], y=4.2)
text(s, ML, 5.75, CW, 0.5,
     "Dürüst not: bu olayda sapma eşzamanlı kanıttır, öngörücü bir öncül değildir.", 12.5, MUTED)
notes(s, "Altın'ın kalbi: birden çok sinyalin baseline'dan saptığını saptıyoruz, sonra alarmları "
         "nedensel önceliğe göre sıralayıp kök nedene iniyoruz. Z-ekseni hatası bir sonuç; gerçek "
         "kök neden hava basıncı. Dürüstlük: sapmanın bu olayda eşzamanlı olduğunu açıkça söyle. "
         "Canlı demo: Tahmin → Aksiyon panosu, bölüm 2.")

# ───────────────────────── 07 · tahmin -> değer ─────────────────────────
s = new("Tahmin · Değer", "ROC / lift katmanından OEE ve €'ya", 7)
stats = [
    ("0,73", "ROC-AUC"),
    ("2,38×", "taban-üstü kazanç"),
    (f"%{DEPLOYED['recall'] * 100:.0f}", "duruş yakalama · held-out"),
    (f"+{DEPLOYED['oee']['delta']['dOEE'] * 100:.2f}".replace(".", ","),
     "puan OEE · e=%35"),
]
sw = (CW - 3 * 0.25) / 4
for i, (v, l) in enumerate(stats):
    x = ML + i * (sw + 0.25)
    box(s, x, 2.55, sw, 1.35, fill=WHITE, line=HAIR, lw=1.0)
    text(s, x, 2.76, sw, 0.55, v, 30, GREEN, True, align=PP_ALIGN.CENTER)
    text(s, x + 0.08, 3.42, sw - 0.16, 0.34, l, 10.5, MUTED, align=PP_ALIGN.CENTER)
bullets(s, [
    f"Gerçek held-out ölçüm: *{DEPLOYED['caught_stops']}/{DEPLOYED['significant_stops']}* duruş "
    f"yakalandı; *{DEPLOYED['episodes']}* alarm döneminin isabeti %{DEPLOYED['episode_precision']*100:.0f}.",
    f"e=%35 varsayımıyla yıllık projeksiyon: *{DEPLOYED['financial']['annualized']['prevented_h']:.0f} saat*; "
    f"fakat 300 €/kontrol varsayımında net *{DEPLOYED['financial']['annualized']['net_eur']:,.0f} €*.",
], y=4.25, gap=8)
box(s, ML, 5.65, CW, 0.72, fill=GREEN_SOFT, line=GREEN, lw=1.2)
text(s, ML + 0.22, 5.83, CW - 0.44, 0.36,
     f"Ekonomik duyarlılık (retrospektif): eşik {ECON['threshold']:.2f} → "
     f"{ECON['episodes']} kontrol → {ECON['annualized_net_eur']:,.0f} €/yıl net. "
     "Canlı eşiğin yerine geçmez.", 12.5, INK, lh=1.1)
notes(s, "Tahmini ilk kez doğrudan OEE ve euroya bağlıyoruz. Recall yakalanabilecek duruş havuzunu, "
         "alarm isabeti ve alarm sayısı ise müdahale maliyetini belirliyor. Varsayılan 300 euro kontrol "
         "maliyetinde dağıtılan eşik ekonomik değil; bunu saklamıyoruz. Retrospektif ekonomik eşik daha "
         "az alarm ve pozitif net değer üretiyor, ama canlı model seçimi olarak sunulmuyor.")

# ───────────────────────── 08 · PLATİN cross ─────────────────────────
s = new("Platin", "Çapraz-makine örüntüleri — null-model ile", 8)
text(s, ML, BODY_Y, CW, 0.4,
     "≥2 makinenin aynı saatte plansız duruşa girmesi — şansın ve vardiya ritminin ötesinde mi?",
     14, INK)
bars = [("Gözlemlenen", 708, GREEN), ("Beklenen (saat-içi ritim)", 564, MUTED),
        ("Beklenen (rastgele)", 333, HAIR)]
maxv = 760
for i, (lab, v, col) in enumerate(bars):
    y = 3.15 + i * 0.62
    text(s, ML, y - 0.02, 3.0, 0.4, lab, 12, INK)
    box(s, ML + 3.1, y, (v / maxv) * 6.2, 0.4, fill=col)
    text(s, ML + 3.2 + (v / maxv) * 6.2, y - 0.01, 1.2, 0.4, f"{v} saat", 12, INK, True)
text(s, ML, 5.2, CW, 0.4,
     "Sonuç:  *z = 5,16, p < 0,001* — senkronizasyon, ortak vardiya programıyla açıklanamıyor.",
     14, INK)
bullets(s, [
    "Eski yöntemin 1 numaralı ‘sistemik’ bulgusu (bağlantı) bir kayıt tekrarıydı; ayıkladık.",
    "Eşleşen küme {1,2,3,9} ortak bir çalışma ritmidir — eşzamanlı arıza yayılımı değil (dürüst okuma).",
], y=5.7, gap=7)
notes(s, "Platin'in en kritik kriteri. Her çapraz-makine iddiamızı bir null-model'den geçirdik. "
         "Eşzamanlı duruşlar gerçek; ama kümeyi abartmadık — ortak bir ritim, eşzamanlı arıza değil. "
         "Canlı demo: Çapraz Makine panosu.")

# ───────────────────────── 09 · PLATİN scenario catalog ─────────────────────────
s = new("Platin", "İncelenebilir senaryo kataloğu", 9)
cols = [("Senaryo / kapsam", 4.15), ("ΔA", 1.05), ("ΔP", 1.05),
        ("ΔOEE", 1.15), ("Saat", 1.2), ("Net €", 1.55)]
x = ML
for label, width in cols:
    box(s, x, 2.35, width, 0.45, fill=INK)
    text(s, x + 0.1, 2.47, width - 0.2, 0.22, label, 10, WHITE, True)
    x += width
for i, row in enumerate(SCENARIOS):
    y = 2.85 + i * 0.68
    fill = GREEN_SOFT if i == 0 else WHITE
    x = ML
    values = [
        (f"{row['id']} · {row['machine']} · {row['scenario']}", 4.15, PP_ALIGN.LEFT),
        (f"{row['delta_A_pp']:.2f}".replace(".", ","), 1.05, PP_ALIGN.RIGHT),
        (f"{row['delta_P_pp']:.2f}".replace(".", ","), 1.05, PP_ALIGN.RIGHT),
        (f"{row['delta_OEE_pp']:.2f}".replace(".", ","), 1.15, PP_ALIGN.RIGHT),
        (f"{max(row['recovered_runtime_h'], row['recovered_schedule_h']):.0f}", 1.2, PP_ALIGN.RIGHT),
        (f"{row['net_eur']:,.0f}", 1.55, PP_ALIGN.RIGHT),
    ]
    for value, width, align in values:
        box(s, x, y, width, 0.6, fill=fill, line=HAIR, lw=0.7)
        text(s, x + 0.1, y + 0.15, width - 0.2, 0.27, value, 10.5,
             GREEN if i == 0 else INK, i == 0, align=align)
        x += width
box(s, ML, 5.78, CW, 0.62, fill=GREEN_SOFT, line=GREEN, lw=1.2)
text(s, ML + 0.22, 5.93, CW - 0.44, 0.3,
     "S2 yalnız sınıflandırmayı değiştirir; S3 üretim sayımı olmadığı için inerttir; "
     "S4 IT aksiyonudur ve makine OEE'sine yazılmaz. Tüm € değerleri VARSAYIM.", 11.5, INK)
notes(s, "Dört senaryo aynı hesap sözleşmesiyle incelenebilir. S1 gerçek OEE ve zaman kazanımıdır. "
         "S2 runtime yaratmadan A'yı yükseltir; bu yüzden finansal faydası yoktur. S3 veri yokluğunda "
         "bilerek inerttir. S4 bağlantı görünürlüğünü geri getirir ama makine OEE'sine yazılmaz.")

# ───────────────────────── 10 · sınırlar ─────────────────────────
s = new("Sınırlar", "Neyi tahmin edemeyeceğimizi de biliyoruz", 10)
bullets(s, [
    "Kondisyon sinyalleri (sıcaklık, yük) boş → tahminin tavanını model değil, *veri* belirler.",
    "Mitsubishi 60 dk’da %53 duruş: tahmin doygun. 30 dk ufkunda zayıf ama mevcut; makineler atılmadı.",
    "Plansız duruşlar tek etiket taşır → ‘ne / ne zaman’ var, ‘neden’ yalnız Makine 1 ve 2’de.",
    "Kalite kaydı yok → Q simüle edilir; üretimsiz günlerde P doğal olarak sıfır.",
], y=BODY_Y, gap=13)
notes(s, "Bunu biz söyleyince güven artar: neyi tahmin edemediğimizi de biliyoruz. Sınırları jüri "
         "bulmadan biz koyuyoruz. Dürüst bir analiz, parlak ama yanlış olandan iyidir.")

# ───────────────────────── 11 · sonuç ─────────────────────────
s = new("Sonuç", "Dürüst, çalışan, uçtan uca", 11)
bullets(s, [
    "Dört seviye de kapsandı; en üst ikisi — Altın ve Platin — güçlü biçimde karşılandı.",
    "Tek akışta: duruş tahmini (Fanuc 2,38×) + kök-neden (kaskad) + sayısal ΔOEE / € önerisi.",
    "Sonraki adım: makine-bazlı olasılık kalibrasyonu, Mitsubishi için kısa-ufuk hedef.",
], y=BODY_Y)
box(s, ML, 5.35, CW, 0.95, fill=INK)
text(s, ML + 0.25, 5.35, CW - 0.5, 0.95,
     "Tahmin eder, açıklar, sayısallaştırır — ve filodaki her makine için neden öyle davrandığımızı biliriz.",
     16, WHITE, True, anchor=MSO_ANCHOR.MIDDLE, lh=1.15)
notes(s, "Kapanış: dürüst, çalışan, uçtan uca bir sistem kurduk. Tahmin eder, açıklar, "
         "sayısallaştırır; ve her scoping kararının veri temelli bir gerekçesi var.")

# ───────────────────────── 12–18 · dashboard guide appendix ─────────────────────────
s = new("Canlı Arayüz Rehberi", "Genel Bakış — tesis ve makine resmi", 12)
picture_fit(s, DASH / "overview_fleet.png", ML, 2.18, 7.95, 4.55)
box(s, 9.05, 2.18, 3.43, 4.55, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.3, 2.48, 2.95, 0.35, "NASIL OKUNUR?", 11, GREEN, True)
bullets(s, [
    "*Üst KPI'lar:* tesis OEE'si, toplam plansız duruş ve Fanuc tahmin lift'i.",
    "*Makine kartları:* renk rejimi, yüzde OEE'yi; alt satır duruş saatini gösterir.",
    "*Sağ panel:* seçilen makinenin A/P/Q ayrışımı ve modelleme rejimi.",
], x=9.3, y=3.0, w=2.9, size=11.5, gap=8)
notes(s, "Genel Bakış ekranında önce üç tesis KPI'sını göster. Sonra makine kartlarının rejim "
         "renklerini açıkla: yeşil Fanuc tahmin hücresi, kehribar Mitsubishi RCA/OEE, gri "
         "telemetrisiz. Bir karta tıklanınca sağda OEE bileşenleri açılır.")

s = new("Canlı Arayüz Rehberi", "Pareto — arıza ile bağlantıyı ayır", 13)
picture_fit(s, DASH / "overview_pareto.png", ML, 2.28, 8.25, 3.25)
box(s, 9.35, 2.28, 3.13, 3.25, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.6, 2.55, 2.65, 0.35, "KARAR MESAJI", 11, GREEN, True)
bullets(s, [
    "*Kırmızı:* giderilebilir makine duruşu.",
    "*Gri:* System Offline; IT/ağ sahipliğinde.",
    "TurboCut büyük kayıp kaynağıdır fakat telemetrisi yoktur.",
], x=9.6, y=3.05, w=2.55, size=11.5, gap=8)
box(s, ML, 5.78, CW, 0.62, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, ML + 0.22, 5.94, CW - 0.44, 0.3,
     "Sunum cümlesi: “Bağlantı kaybını makine arızası sayarsak yanlış ekibe ve yanlış yatırıma gideriz.”",
     11.5, INK)
notes(s, "Pareto'da kırmızı ve gri ayrımını vurgula. System Offline makine arızası değildir. "
         "Bu ekran aksiyon sahibini belirler: bakım mı, IT/ağ mı?")

s = new("Canlı Arayüz Rehberi", "Tahmin — risk zaman çizgisi ve gerçek duruşlar", 14)
picture_fit(s, DASH / "predict_risk.png", ML, 2.15, 8.25, 4.78)
box(s, 9.35, 2.15, 3.13, 4.78, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.6, 2.45, 2.65, 0.35, "NE GÖSTERİYOR?", 11, GREEN, True)
bullets(s, [
    "*Yeşil çizgi:* modelin duruş riski.",
    "*Kırmızı kesik:* denetlenen alarm eşiği 0,1778.",
    "*Kırmızı çarpılar:* gerçekleşen ≥15 dk plansız duruşlar.",
    "Alt kutular, yüksek-risk dönemlerini inceleme sırasına koyar.",
], x=9.6, y=2.95, w=2.55, size=11.3, gap=7)
notes(s, "Bu ekran modelin tutulmamış gelecekteki riskini gösterir. Yeşil risk çizgisini, sabit "
         "alarm eşiğini ve gerçek duruş işaretlerini anlat. Alt dönemlerden biri seçilince RCA "
         "ve What-If akışı aynı sayfada devam eder.")

s = new("Canlı Arayüz Rehberi", "Kök-Neden — alarm kaskadından aksiyona", 15)
picture_fit(s, DASH / "predict_rca.png", ML, 2.28, 8.25, 3.05)
box(s, 9.35, 2.28, 3.13, 3.05, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.6, 2.55, 2.65, 0.35, "ALTIN AKIŞI", 11, GREEN, True)
bullets(s, [
    "Kaskad nedensel önceliğe göre sıralanır.",
    "*AIR PRESSURE* kök neden; Z-axis alarmı sonuçtur.",
    "Telemetri ve hipotezler önerilen aksiyonu destekler.",
], x=9.6, y=3.0, w=2.55, size=11.3, gap=7)
box(s, ML, 5.68, CW, 0.72, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, ML + 0.22, 5.85, CW - 0.44, 0.38,
     "Sunum cümlesi: “Model yalnız alarm vermiyor; alarmı kanıt, hipotez ve uygulanabilir aksiyona çeviriyor.”",
     11.5, INK)
notes(s, "Makine 1 örneğinde alarm sırası ile neden sırasının aynı şey olmadığını anlat. "
         "Hava basıncı kök neden, Z ekseni geri dönüş alarmı sonuçtur. Sağdaki hipotezler "
         "olasılık ve önerilen aksiyonla birlikte gösterilir.")

s = new("Canlı Arayüz Rehberi", "Kestirimci bakım değeri ve What-If", 16)
picture_fit(s, DASH / "predict_value.png", ML, 2.15, 8.0, 2.18)
picture_fit(s, DASH / "predict_whatif.png", ML, 4.52, 8.0, 2.25)
box(s, 9.1, 2.15, 3.38, 4.62, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.35, 2.45, 2.9, 0.35, "İKİ FARKLI HESAP", 11, GREEN, True)
bullets(s, [
    "*Üst panel:* yalnız modelin yakaladığı duruşlara atfedilen OEE/€.",
    "Etkililik ve maliyetler varsayımdır; recall/precision held-out ölçümdür.",
    "*Alt panel:* kullanıcı seçtiği duruş azaltımını genel What-If motorunda simüle eder.",
    "Negatif ROI saklanmaz; alarm maliyeti her gerçek/yanlış alarm için yazılır.",
], x=9.35, y=2.95, w=2.85, size=11.15, gap=7)
notes(s, "Üst kart model-atfedilebilir değerdir: kaçırılan duruşlara değer yazılmaz ve her alarm "
         "kontrol maliyeti doğurur. Alt kart genel What-If'tir; kullanıcı yüzdeyi ve finansal "
         "varsayımları değiştirerek A/P/Q/OEE etkisini anında görür.")

s = new("Canlı Arayüz Rehberi", "Senaryo kataloğu — sonuçlar denetlenebilir", 17)
picture_fit(s, DASH / "predict_scenarios.png", ML, 2.25, 8.45, 2.48)
box(s, 9.45, 2.25, 3.03, 2.48, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.7, 2.52, 2.55, 0.35, "TABLO NE İŞE YARAR?", 11, GREEN, True)
bullets(s, [
    "ΔA / ΔP / ΔQ / ΔOEE birlikte görünür.",
    "Runtime ve schedule kazanımı ayrıdır.",
    "Sütun başlığına tıklayarak sıralanır.",
], x=9.7, y=2.95, w=2.45, size=11.2, gap=7)
box(s, ML, 5.18, CW, 1.18, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, ML + 0.22, 5.42, CW - 0.44, 0.72,
     "S1 gerçek zaman/OEE kazanımıdır. S2 yalnız sınıflandırma etkisini, S3 üretim verisi yokluğundaki "
     "inert kolu, S4 ise IT sahipliğindeki schedule görünürlüğünü gösterir.", 12, INK, lh=1.16)
notes(s, "Katalogdaki dört satırı kısa yorumla. Özellikle S2'nin runtime yaratmadığını, S3'ün "
         "ProductSum sıfır olduğu için inert kaldığını ve S4'ün makine OEE'sine yazılmadığını söyle.")

s = new("Canlı Arayüz Rehberi", "Çapraz Makine — iddiayı null-model ile sınar", 18)
picture_fit(s, DASH / "cross_machine.png", ML, 2.12, 8.1, 4.82)
box(s, 9.05, 2.12, 3.43, 4.82, fill=GREEN_SOFT, line=GREEN, lw=1.1)
text(s, 9.3, 2.42, 2.95, 0.35, "PLATİN EKRANI", 11, GREEN, True)
bullets(s, [
    "Eski connectivity sonucu kayıt tekrarı olarak teşhis edilir.",
    "708 eşzamanlı duruş, vardiya ritmini koruyan null'ın üstündedir.",
    "Rejim haritası hangi makinelerin birlikte modellenebileceğini gösterir.",
    "Küme seviyesi güçlü olsa da türev korelasyonu ≈0: akut yayılım iddia edilmez.",
], x=9.3, y=2.92, w=2.9, size=11.05, gap=7)
notes(s, "Bu ekranı soru gelirse aç. Eski yöntemin totolojisini, null-model destekli "
         "senkronizasyonu, veri rejimlerini ve kümenin dürüst yorumunu tek ekranda gösterir.")

prs.save(OUT)
print(f"kaydedildi: {OUT}  ({len(prs.slides)} slayt)")
